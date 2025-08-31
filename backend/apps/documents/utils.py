import PyPDF2
import docx
from pptx import Presentation
import magic
import os
from django.conf import settings


def get_file_type(file_path):
    
    """Determine le type de fichier"""
    
    mime        = magic.Magic(mime=True)
    file_type   = mime.from_file(file_path)
    
    if 'pdf' in file_type:
        return 'pdf'
    elif 'wordprocessingml' in file_type or 'msword' in file_type:
        return 'docx'
    elif 'presentationml' in file_type or 'presentation' in file_type:
        return 'pptx'
    else:
        return 'unknown'
    


def extract_text_from_pdf(file_path):
    
    """Extrait le texte d'un fichier pdf"""
    
    try:
        with open(file_path, 'rb') as file:
            reader  = PyPDF2.PdfReader(file)
            text    = ""
            for page in reader.pages:
                text += page.extract_text() + '\n'
            return text.strip()
    except Exception as e:
        print(f"Erreur lors de l'extraction PDF {e}")
        return ""
    


def extract_text_from_docx(file_path):
    
    """Extrait le texte d'un fichier docx"""
    
    try:
        doc     = docx.Document(file_path)
        text    = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + '\n'
        return text.strip()
    except Exception as e:
        print(f"Erreur lors de l'extractin DOCX {e}")
        return ""


def extract_text_from_pptx(file_path):
    
    """Extrait le text d'un fichier PPTX"""
    
    try:
        prs     = Presentation(file_path)
        text    = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + '\n'
        return text.strip()
    except Exception as e:
        print(f"Erreur lors de l'extraction du PPTX {e}")
        return ""
    

def extract_text_from_document(file_path, file_type):
    
    """Extrait le texte selon le type de fichier"""
    
    if file_type == 'pdf':
        return extract_text_from_pdf(file_path)
    elif file_type == 'docx':
        return extract_text_from_docx(file_path)
    elif file_type == 'pptx':
        return extract_text_from_pptx(file_path)
    else:
        return ""
    

def classify_document(text):
    
    """Classifie un document en se basant sur son contenu"""
    
    text_lower = text.lower()
    categories = settings.DOCUMENT_CATEGORIES
    
    category_scores = {}
    
    for category, keywords in categories.items():
        if category == "autres":
            continue
        score = 0
        found_keywords = []
        
        for keyword in keywords:
            count = text_lower.count(keyword)
            if count > 0:
                score += count
                found_keywords.append(keyword)
        if score > 0:
            category_scores[category] = {
                'score': score,
                'keywords': found_keywords
            }
    if category_scores:
        # retourner la category avec le score le plus eleve
        best_category = max(category_scores.items(), key=lambda x: x[1]['score'])
        return best_category[0], best_category[1]['keywords']
    else:
        return 'autres', []